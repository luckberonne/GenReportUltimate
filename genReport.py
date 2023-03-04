from flask import Flask, render_template, request, redirect, url_for
import sqlite3
import csv
from pptx import Presentation
from pptx.util import Inches
import os
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication


app = Flask(__name__)


class PresentationCreator:
    def __init__(self, db_conn):
        self.db_conn = db_conn

    def read_data(self, table_name):
        c = self.db_conn.cursor()
        c.execute(f"SELECT name, age, email FROM {table_name}")
        data = c.fetchall()
        return data

    def write_data_to_csv(self, data, filename):
        with open(filename, "w", newline="") as csvfile:
            writer = csv.writer(csvfile)
            writer.writerow(["Name", "Age", "Email"])
            writer.writerows(data)

    def create_title_slide(self, prs, title):
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        return slide

    def create_data_slides(self, prs, data):
        for row in data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = row[0]
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = f"Age: {row[1]}\nEmail: {row[2]}"
            img_path = f"{row[0]}.jpg"
            if os.path.exists(img_path):
                pic = slide.shapes.add_picture(img_path, Inches(4), Inches(3))

    def create_presentation(self, data, title):
        prs = Presentation()
        self.create_title_slide(prs, title)
        self.create_data_slides(prs, data)
        return prs

    def send_email(self, sender_email, sender_password, recipient_email, subject, body, files=None):
        msg = MIMEMultipart()
        msg["From"] = sender_email
        msg["To"] = recipient_email
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain"))
        for f in files or []:
            with open(f, "rb") as attachment:
                part = MIMEApplication(attachment.read(), Name=os.path.basename(f))
                part["Content-Disposition"] = f'attachment; filename="{os.path.basename(f)}"'
                msg.attach(part)
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()
            server.login(sender_email, sender_password)
            text = msg.as_string()
            server.sendmail(sender_email, recipient_email, text)

    def insert_data_from_csv(self, filename, table_name):
        with open(filename, newline='') as csvfile:
            reader = csv.reader(csvfile)
            next(reader)  # skip header row
            c = self.db_conn.cursor()
            for row in reader:
                c.execute(f'INSERT INTO {table_name} (name, age, email) VALUES (?, ?, ?)', row)
        self.db_conn.commit()

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        file = request.files['file']
        if file.filename == '':
            return redirect(request.url)
        if file:
            file.save(file.filename)
            conn = sqlite3.connect('database.db')
            creator = PresentationCreator(conn)
            creator.insert_data_from_csv(file.filename, 'people')
            conn.close()
            data = creator.read_data('people')
            prs = creator.create_presentation(data, 'People Information')
            prs.save('people_info.pptx')
            sender_email = 'sender@gmail.com'
            sender_password = 'password'
            recipient_email = 'recipient@gmail.com'
            subject = 'People Information'
            body = 'Please find attached the people information file.'
            files = ['people_info.pptx']
            creator.send_email(sender_email, sender_password, recipient_email, subject, body, files)
            os.remove(file.filename)
            os.remove('people_info.pptx')
            return redirect(url_for('index'))
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)

# @app.route('/')
# def home():
#     name = "Worldhh"
#     return render_template("index.html", name=name)

# # Run the app
# if __name__ == '__main__':
#     app.run()